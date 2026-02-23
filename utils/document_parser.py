import os
from PyPDF2 import PdfReader
from pptx import Presentation


def parse_pdf(file) -> str:
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text.strip()


def parse_pptx(file_or_path) -> str:
    prs = Presentation(file_or_path)
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text.strip()


def parse_document(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    if name.endswith(".pdf"):
        return parse_pdf(uploaded_file)
    elif name.endswith((".pptx", ".ppt")):
        return parse_pptx(uploaded_file)
    elif name.endswith((".txt", ".md", ".csv")):
        return uploaded_file.read().decode("utf-8")
    else:
        return uploaded_file.read().decode("utf-8", errors="ignore")


def parse_document_from_path(file_path: str) -> str:
    name = os.path.basename(file_path).lower()
    if name.endswith(".pdf"):
        with open(file_path, "rb") as f:
            return parse_pdf(f)
    elif name.endswith((".pptx", ".ppt")):
        return parse_pptx(file_path)
    elif name.endswith((".txt", ".md", ".csv")):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
